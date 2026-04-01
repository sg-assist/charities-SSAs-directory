#!/usr/bin/env python3
"""Generate PPTX presentation for UNFPA Partnership Catalyst client briefing."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
LIGHT_BLUE = RGBColor(0x00, 0x9E, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, height=Inches(1.2), color=DARK_BLUE):
    shp = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_accent_bar(slide, top=Inches(1.2), height=Inches(0.08)):
    shp = slide.shapes.add_shape(1, Inches(0), top, prs.slide_width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = LIGHT_BLUE
    shp.line.fill.background()

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullets(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY, bold_prefix=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.space_before = Pt(4)
        if isinstance(item, tuple):
            run = p.add_run()
            run.text = item[0]
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = DARK_BLUE
            run.font.name = 'Calibri'
            run2 = p.add_run()
            run2.text = item[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = 'Calibri'
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = 'Calibri'
    return tf

def slide_title(slide, title, subtitle=None):
    add_title_bar(slide)
    add_accent_bar(slide)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8), title, font_size=30, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.5), subtitle, font_size=16, color=RGBColor(0xBB, 0xDD, 0xFF))

# ── SLIDE 1: Title ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK_BLUE)
bar = s.shapes.add_shape(1, Inches(0), Inches(5.5), prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = LIGHT_BLUE
bar.line.fill.background()
add_text_box(s, Inches(1), Inches(1.5), Inches(11), Inches(2), 'Designing Innovative PPPs for\nClimate & Humanitarian Resilience in Asia', font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
add_text_box(s, Inches(1), Inches(3.8), Inches(11), Inches(0.6), 'UNFPA Partnership Catalyst', font_size=24, color=LIGHT_BLUE, alignment=PP_ALIGN.LEFT)
add_text_box(s, Inches(1), Inches(4.6), Inches(11), Inches(1), 'LKYSPP Policy Innovation Lab  |  National University of Singapore\nProfessor Mancini  |  April 2026', font_size=16, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.LEFT)

# ── SLIDE 2: HMW Question ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_title_bar(s, Inches(0.6))
add_accent_bar(s, Inches(0.6))
add_text_box(s, Inches(0.8), Inches(0.1), Inches(11.5), Inches(0.5), 'The Design Challenge', font_size=24, bold=True, color=WHITE)
shp = s.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5))
shp.fill.solid()
shp.fill.fore_color.rgb = RGBColor(0xF0, 0xF6, 0xFC)
shp.line.color.rgb = LIGHT_BLUE
shp.line.width = Pt(2)
add_text_box(s, Inches(1.8), Inches(1.6), Inches(9.7), Inches(0.5), 'HOW MIGHT WE...', font_size=18, bold=True, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(2), Inches(2.3), Inches(9.3), Inches(3.5), 'Design innovative public-private partnerships that enable fragmented climate and humanitarian finance to flow efficiently toward inclusive, community-led resilience initiatives across Asia?', font_size=28, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# ── SLIDE 3: The Challenge ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, 'The Challenge: Fragmented Finance')
add_bullets(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5), [
    ('$350B+ ', 'in annual climate finance flows globally, but less than 10% reaches local communities in Asia'),
    ('Humanitarian funding ', 'operates in parallel silos: bilateral donors, multilateral funds, private philanthropy, DFIs'),
    ('SRHR-climate nexus ', 'is severely underfunded despite clear evidence linking reproductive health to climate resilience'),
    ('Community-led initiatives ', 'struggle to access large-scale funding due to scale mismatch, reporting requirements, and intermediary costs'),
    ('Result: ', 'billions in available capital, but fragmented across hundreds of mechanisms with no common intelligence layer'),
], font_size=18)

# ── SLIDE 4: UNFPA Position ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "UNFPA's Unique Position")
add_bullets(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), [
    ('150+ countries ', 'with operational presence and government partnerships'),
    ('$1.1 billion/year ', 'in programme delivery across SRHR, maternal health, and humanitarian response'),
    ('Three Transformative Results: ', 'zero preventable maternal deaths, zero unmet need for family planning, zero GBV and harmful practices'),
    ('World\'s largest ', 'contraceptive procurer with established global supply chains'),
    ('Mandate from ICPD 1994: ', 'rights-based approach to population and development'),
], font_size=17)
add_text_box(s, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5), 'UNFPA sits at the intersection of:\n\n'
    + '   Sexual & reproductive health\n'
    + '   + Climate adaptation\n'
    + '   + Humanitarian response\n'
    + '   + Community resilience\n\n'
    + 'This intersection is exactly where\ninnovative PPPs are needed most.',
    font_size=17, color=DARK_BLUE)

print("Part 1 done - slides 1-4")
# ── SLIDE 5: The Knowledge Gap ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, 'The Knowledge Gap')
add_bullets(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5), [
    'UNFPA staff preparing for funder meetings must manually research across hundreds of PDFs, reports, and evaluations',
    'Institutional knowledge about programme areas, evidence quality, and contested topics is scattered and hard to synthesise',
    'Real-time information about funder priorities, market conditions, and policy changes requires separate research',
    'No single tool combines institutional intelligence with current market data for partnership preparation',
    'Staff need to frame UNFPA programmes in the language of funders: climate finance, blended finance, impact investment',
], font_size=18)

# ── SLIDE 6: Our Solution ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, 'Our Solution: The Partnership Catalyst')
# Two columns
add_text_box(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5), 'Knowledge Base', font_size=22, bold=True, color=DARK_BLUE)
add_bullets(s, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4), [
    '30+ deep-research analytical documents',
    'Covering 6 thematic blocks across UNFPA',
    '100+ ingested PDF source reports',
    'Browsable, searchable, cross-referenced',
    'Honest assessment of contested areas',
    'Directly addresses the HMW question through the Resilience & Partnerships block',
], font_size=16)
add_text_box(s, Inches(7), Inches(1.6), Inches(5.5), Inches(0.5), 'AI Chat Engine', font_size=22, bold=True, color=DARK_BLUE)
add_bullets(s, Inches(7), Inches(2.2), Inches(5.5), Inches(4), [
    'Agentic Claude Sonnet 4 with tool use',
    'Semantic search across entire knowledge base',
    'Real-time web search for current information',
    'Extended thinking for complex analysis',
    'Source attribution with clickable links',
    'Generates pitches, briefings, talking points',
], font_size=16)
# Divider line
shp = s.shapes.add_shape(1, Inches(6.5), Inches(1.6), Inches(0.04), Inches(4.5))
shp.fill.solid()
shp.fill.fore_color.rgb = LIGHT_BLUE
shp.line.fill.background()
add_text_box(s, Inches(2), Inches(6.5), Inches(9), Inches(0.5), 'Live at: unfpa-lkyspp-otg.vercel.app', font_size=14, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

# ── SLIDE 7: Knowledge Base Overview ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, 'The Knowledge Base: 30+ Analytical Documents')
add_bullets(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.5), [
    'Each document is 3,000-14,000 words of original research and analysis',
    'Written to answer specific questions practitioners and decision-makers actually ask',
    'Not summaries of existing reports \u2014 original analytical synthesis with source attribution',
], font_size=17)
# Table
from pptx.util import Inches, Pt
tbl = s.shapes.add_table(7, 4, Inches(0.8), Inches(3.2), Inches(11.5), Inches(3.8)).table
tbl.columns[0].width = Inches(2.5)
tbl.columns[1].width = Inches(1)
tbl.columns[2].width = Inches(1)
tbl.columns[3].width = Inches(7)
headers = ['Block', 'Code', 'Docs', 'Coverage']
rows = [
    ['Orientation', 'O', '8', 'What UNFPA/PMNCH are, mandate, structure, terminology'],
    ['Programme Work', 'W', '10', 'Maternal health, family planning, GBV, midwifery, FGM, child marriage'],
    ['Data & Evidence', 'D', '4', 'SOWP reports, census, demographic dividend, results methodology'],
    ['Contested Areas', 'C', '5', 'US defunding, abortion, CSE, disputed results, China programme'],
    ['PMNCH', 'PMNCH', '5', 'Partnership mandate, advocacy, accountability, overlap with UNFPA'],
    ['Resilience & PPPs', 'R', '4', 'PPP models, climate-SRHR, Singapore ecosystem, community resilience'],
]
for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            p.font.name = 'Calibri'
        if ri % 2 == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

print("Part 2 done - slides 5-7")
# ── SLIDE 8: Resilience & Partnerships Block ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, 'Resilience & Partnerships: Addressing the HMW Directly', 'Block R \u2014 4 documents purpose-built for the climate-resilience PPP challenge')
add_bullets(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), [
    ('UNFPA-R-01: PPP Models ', '\u2014 Taxonomy of partnership types (commodity procurement, blended finance, platforms, impact bonds). Evidence on what works and what fails. Design principles for community accountability.'),
    ('UNFPA-R-02: Climate & SRHR ', '\u2014 The evidence linking climate change to reproductive health outcomes in Asia. How to frame SRHR within climate funding mechanisms.'),
    ('UNFPA-R-03: Singapore Ecosystem ', '\u2014 Singapore\'s financial infrastructure for resilience: family offices, DFIs, sovereign wealth, impact investors. How to engage each.'),
    ('UNFPA-R-04: Community Resilience ', '\u2014 Co-design principles, intergenerational models, community ownership. How to ensure finance reaches community-led initiatives.'),
], font_size=15)
# Right side callout
shp = s.shapes.add_shape(1, Inches(7), Inches(1.8), Inches(5.5), Inches(4))
shp.fill.solid()
shp.fill.fore_color.rgb = RGBColor(0xF0, 0xF6, 0xFC)
shp.line.color.rgb = LIGHT_BLUE
shp.line.width = Pt(1)
add_text_box(s, Inches(7.3), Inches(2), Inches(5), Inches(3.5),
    'These four documents directly answer\nthe components of the HMW:\n\n'
    + '\u2714 Innovative PPP design (R-01)\n'
    + '\u2714 Climate-humanitarian finance (R-02)\n'
    + '\u2714 Efficient flow of capital (R-03)\n'
    + '\u2714 Community-led resilience (R-04)',
    font_size=16, color=DARK_BLUE)

# ── SLIDE 9: AI Chat Engine ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, 'The AI Chat Engine: Agentic Intelligence')
add_bullets(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), [
    ('Claude Sonnet 4 ', '\u2014 Anthropic\'s most capable model with extended thinking (8K token budget) for deep reasoning'),
    ('Semantic Search (RAG) ', '\u2014 Vector similarity search across all knowledge base chunks using pgvector. Finds the most relevant content for any query.'),
    ('Real-Time Web Search ', '\u2014 Anthropic\'s built-in web search tool fetches current funder information, market data, and policy changes.'),
    ('Agentic Loop ', '\u2014 Claude decides which tools to use and when. Can search knowledge base, then web, then knowledge base again \u2014 up to 6 rounds.'),
    ('Wave Generation ', '\u2014 Automatically continues long responses across up to 4 waves, preventing truncated answers.'),
], font_size=16)
add_text_box(s, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5),
    'What This Means in Practice:\n\n'
    + 'A user asks: "Help me pitch our\n'
    + 'climate-SRHR programme to a Singapore\n'
    + 'family office."\n\n'
    + 'Claude:\n'
    + '1. Thinks through the question\n'
    + '2. Searches KB for SRHR + PPP docs\n'
    + '3. Searches web for family office interests\n'
    + '4. Synthesises a tailored pitch\n'
    + '5. Cites all sources used',
    font_size=15, color=DARK_BLUE)

# ── SLIDE 10: How It Works ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, 'How It Works: From Question to Answer')
# Flow boxes
flow_items = [
    ('User Question', 'Natural language query\nabout UNFPA, funders,\nor partnerships'),
    ('Extended\nThinking', 'Claude reasons through\nthe question with\n8K token budget'),
    ('Knowledge\nBase Search', 'Semantic vector search\nacross 30+ documents\nand 100+ PDFs'),
    ('Web Search', 'Real-time search for\ncurrent funder data\nand market intelligence'),
    ('Response\nGeneration', 'Source-attributed\nanalysis with citations\nand recommendations'),
]
box_w = Inches(2.2)
gap = Inches(0.15)
start_x = Inches(0.5)
for i, (title, desc) in enumerate(flow_items):
    x = start_x + i * (box_w + gap)
    shp = s.shapes.add_shape(1, x, Inches(2), box_w, Inches(3.5))
    shp.fill.solid()
    shp.fill.fore_color.rgb = DARK_BLUE if i == 0 or i == 4 else RGBColor(0xF0, 0xF6, 0xFC)
    shp.line.color.rgb = LIGHT_BLUE
    shp.line.width = Pt(1)
    tc = WHITE if i == 0 or i == 4 else DARK_BLUE
    add_text_box(s, x + Inches(0.15), Inches(2.2), box_w - Inches(0.3), Inches(1), title, font_size=16, bold=True, color=tc, alignment=PP_ALIGN.CENTER)
    dc = RGBColor(0xCC, 0xDD, 0xEE) if i == 0 or i == 4 else DARK_GRAY
    add_text_box(s, x + Inches(0.15), Inches(3.2), box_w - Inches(0.3), Inches(2), desc, font_size=13, color=dc, alignment=PP_ALIGN.CENTER)
# Arrows
for i in range(4):
    x = start_x + (i + 1) * (box_w + gap) - gap + Inches(0.02)
    add_text_box(s, x, Inches(3.2), gap, Inches(0.5), '\u2192', font_size=24, bold=True, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(s, Inches(0.8), Inches(6), Inches(11.5), Inches(0.5),
    'The agentic loop runs up to 6 rounds \u2014 Claude decides when it has enough information to answer comprehensively.',
    font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

print("Part 3 done - slides 8-10")
# ── SLIDE 11: Use Case - Funder Pitch ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, 'Use Case: Preparing a Funder Pitch', 'Singapore family office interested in climate adaptation')
# Scenario box
shp = s.shapes.add_shape(1, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5))
shp.fill.solid()
shp.fill.fore_color.rgb = RGBColor(0xF0, 0xF6, 0xFC)
shp.line.color.rgb = LIGHT_BLUE
shp.line.width = Pt(1)
add_text_box(s, Inches(1), Inches(1.8), Inches(5), Inches(0.4), 'The Scenario', font_size=18, bold=True, color=DARK_BLUE)
add_text_box(s, Inches(1), Inches(2.3), Inches(5), Inches(4),
    'A UNFPA country representative has a meeting\n'
    + 'with a Singapore-based family office that focuses\n'
    + 'on climate adaptation in Southeast Asia.\n\n'
    + 'They ask the Partnership Catalyst:\n\n'
    + '"Help me pitch our climate-SRHR programme\n'
    + 'for a family office interested in climate\n'
    + 'adaptation in Southeast Asia."',
    font_size=15, color=DARK_GRAY)
# Response box
shp2 = s.shapes.add_shape(1, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5))
shp2.fill.solid()
shp2.fill.fore_color.rgb = DARK_BLUE
shp2.line.fill.background()
add_text_box(s, Inches(7), Inches(1.8), Inches(5.2), Inches(0.4), 'The Catalyst Responds', font_size=18, bold=True, color=LIGHT_BLUE)
add_text_box(s, Inches(7), Inches(2.3), Inches(5.2), Inches(4),
    '\u2714 Searches KB for UNFPA climate-SRHR evidence\n'
    + '\u2714 Searches KB for Singapore financial ecosystem\n'
    + '\u2714 Searches web for family office\'s recent interests\n'
    + '\u2714 Generates tailored pitch with:\n'
    + '    \u2022 Programme impact data from UNFPA-R-02\n'
    + '    \u2022 Singapore engagement strategy from R-03\n'
    + '    \u2022 Community ownership model from R-04\n'
    + '    \u2022 PPP structure options from R-01\n'
    + '\u2714 Cites all knowledge base sources used',
    font_size=15, color=WHITE)

# ── SLIDE 12: Use Case - Partnership Structuring ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, 'Use Case: Partnership Structuring', 'Comparing financing models for a resilience programme')
add_text_box(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.8),
    'A partnership team asks: "Compare blended finance, impact bonds, and South-South cooperation models for a community resilience programme in the Philippines."',
    font_size=17, color=DARK_GRAY)
# Three columns
cols = [
    ('Blended Finance', [
        'Combines public + private capital',
        'De-risks private investment',
        'Strong evidence at platform level',
        'Scale mismatch with community initiatives',
        'Best for: Large-scale infrastructure',
    ]),
    ('Impact Bonds', [
        'Outcomes-based financing',
        'Transfers risk to investors',
        'Complex to structure',
        'High transaction costs for small projects',
        'Best for: Measurable health outcomes',
    ]),
    ('South-South Cooperation', [
        'Peer learning between countries',
        'Lower power asymmetry',
        'Harder to mobilise large capital',
        'Strong community ownership potential',
        'Best for: Knowledge + capacity building',
    ]),
]
for i, (title, items) in enumerate(cols):
    x = Inches(0.8 + i * 4.1)
    add_text_box(s, x, Inches(2.8), Inches(3.8), Inches(0.5), title, font_size=18, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    bar = s.shapes.add_shape(1, x, Inches(3.3), Inches(3.8), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_BLUE
    bar.line.fill.background()
    add_bullets(s, x, Inches(3.5), Inches(3.8), Inches(3.5), items, font_size=14)

add_text_box(s, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5),
    'The Catalyst draws from UNFPA-R-01 (PPP Models) to provide evidence-based comparison with practical recommendations.',
    font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ── SLIDE 13: Connecting Finance to Community ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, 'Connecting Finance to Community-Led Resilience')
add_text_box(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1),
    'The Partnership Catalyst bridges the gap between fragmented global finance and community-led initiatives by providing:',
    font_size=18, color=DARK_GRAY)
# Four boxes
bridge_items = [
    ('Evidence Translation', 'Converts UNFPA programme\nevidence into funder-\nfriendly language and\nimpact frameworks'),
    ('Funder Intelligence', 'Real-time web search\nfor funder priorities,\nrecent investments, and\npartnership preferences'),
    ('Partnership Design', 'PPP models, blended\nfinance structures, and\ncommunity accountability\nframeworks from Block R'),
    ('Community Framing', 'Ensures community\nownership, co-design,\nand intergenerational\nequity in every pitch'),
]
for i, (title, desc) in enumerate(bridge_items):
    x = Inches(0.5 + i * 3.15)
    shp = s.shapes.add_shape(1, x, Inches(2.8), Inches(2.9), Inches(3.5))
    shp.fill.solid()
    shp.fill.fore_color.rgb = DARK_BLUE if i % 2 == 0 else RGBColor(0xF0, 0xF6, 0xFC)
    shp.line.color.rgb = LIGHT_BLUE
    shp.line.width = Pt(1)
    tc = WHITE if i % 2 == 0 else DARK_BLUE
    dc = RGBColor(0xCC, 0xDD, 0xEE) if i % 2 == 0 else DARK_GRAY
    add_text_box(s, x + Inches(0.2), Inches(3), Inches(2.5), Inches(0.6), title, font_size=16, bold=True, color=tc, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x + Inches(0.2), Inches(3.7), Inches(2.5), Inches(2.2), desc, font_size=14, color=dc, alignment=PP_ALIGN.CENTER)

print("Part 4 done - slides 11-13")
# ── SLIDE 14: Technology Architecture ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, 'Technology Architecture')
tbl = s.shapes.add_table(9, 3, Inches(1.5), Inches(1.8), Inches(10), Inches(5)).table
tbl.columns[0].width = Inches(2.5)
tbl.columns[1].width = Inches(3.5)
tbl.columns[2].width = Inches(4)
tech_headers = ['Layer', 'Technology', 'Purpose']
tech_rows = [
    ['Frontend', 'Next.js 16 + React 19', 'Server-rendered pages, interactive chat UI'],
    ['AI Engine', 'Claude Sonnet 4 (Anthropic)', 'Agentic chat with extended thinking + tool use'],
    ['Embeddings', 'OpenAI text-embedding-3-small', '1536-dimension vectors for semantic search'],
    ['Database', 'PostgreSQL + pgvector', 'Document storage + vector similarity search'],
    ['Web Search', 'Anthropic web_search tool', 'Real-time funder and market intelligence'],
    ['Rate Limiting', 'Upstash Redis', '20 queries/day per user'],
    ['Hosting', 'Vercel', 'Serverless deployment, edge functions'],
    ['Language', 'TypeScript', 'Full-stack type safety'],
]
for ci, h in enumerate(tech_headers):
    cell = tbl.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE
for ri, row in enumerate(tech_rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_GRAY
            p.font.name = 'Calibri'
        if ri % 2 == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

# ── SLIDE 15: What's Next ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, "What's Next: Expansion Roadmap")
add_bullets(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), [
    ('Expand Knowledge Base ', '\u2014 Add country-specific documents for priority Asia-Pacific markets. Ingest new evaluation reports and funder databases.'),
    ('Funder Profiles ', '\u2014 Build structured profiles of key funders (DFIs, family offices, sovereign wealth funds) with investment theses and engagement histories.'),
    ('Partnership Templates ', '\u2014 Pre-built frameworks for common PPP structures: blended finance vehicles, impact bonds, co-investment platforms.'),
    ('Multi-Language Support ', '\u2014 Enable queries and responses in languages used across UNFPA Asia-Pacific offices.'),
    ('Usage Analytics ', '\u2014 Track which topics are most queried to identify knowledge gaps and guide new document creation.'),
], font_size=16)
shp = s.shapes.add_shape(1, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5))
shp.fill.solid()
shp.fill.fore_color.rgb = DARK_BLUE
shp.line.fill.background()
add_text_box(s, Inches(7.3), Inches(2), Inches(5), Inches(4),
    'The Vision\n\n'
    + 'Every UNFPA staff member preparing\n'
    + 'for a funder conversation has instant\n'
    + 'access to:\n\n'
    + '\u2022 Deep institutional intelligence\n'
    + '\u2022 Current funder data\n'
    + '\u2022 Evidence-based talking points\n'
    + '\u2022 Partnership design frameworks\n'
    + '\u2022 Community accountability models\n\n'
    + 'All in one conversation.',
    font_size=16, color=WHITE)

# ── SLIDE 16: Team ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
slide_title(s, 'Team & Attribution')
add_text_box(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.5), 'LKYSPP Student Group', font_size=20, bold=True, color=DARK_BLUE)
add_bullets(s, Inches(0.8), Inches(2.4), Inches(5.5), Inches(2.5), [
    'Rani Opula Rajan',
    'Prachi Sharma',
    'Abhishek Tiwari',
    'Preeti Patil',
], font_size=16)
add_text_box(s, Inches(0.8), Inches(4.6), Inches(5.5), Inches(0.5), 'Application', font_size=20, bold=True, color=DARK_BLUE)
add_text_box(s, Inches(0.8), Inches(5.2), Inches(5.5), Inches(1),
    'Designed and built by Haojun See (MPP 2021)\nOn The Ground AI',
    font_size=16, color=DARK_GRAY)
add_text_box(s, Inches(7), Inches(1.8), Inches(5.5), Inches(0.5), 'Academic Context', font_size=20, bold=True, color=DARK_BLUE)
add_text_box(s, Inches(7), Inches(2.4), Inches(5.5), Inches(3),
    'Policy Innovation Lab\n'
    + 'Professor Mancini\n\n'
    + 'Lee Kuan Yew School of Public Policy\n'
    + 'National University of Singapore\n\n'
    + 'Client: UNFPA\n'
    + 'Challenge B: PPP for Climate and\n'
    + 'Humanitarian Resilience',
    font_size=16, color=DARK_GRAY)

# ── SLIDE 17: Thank You ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK_BLUE)
bar = s.shapes.add_shape(1, Inches(0), Inches(5.5), prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = LIGHT_BLUE
bar.line.fill.background()
add_text_box(s, Inches(1), Inches(2), Inches(11), Inches(1.5), 'Thank You', font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1), Inches(3.5), Inches(11), Inches(1),
    'unfpa-lkyspp-otg.vercel.app', font_size=22, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1), Inches(4.5), Inches(11), Inches(1),
    'UNFPA@ontheground.agency', font_size=16, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)
add_text_box(s, Inches(1), Inches(5.8), Inches(11), Inches(0.8),
    'LKYSPP Policy Innovation Lab  |  Professor Mancini  |  National University of Singapore',
    font_size=14, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

# ── SAVE ──
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'docs', 'deliverables',
                        'UNFPA-Partnership-Catalyst-Briefing.pptx')
prs.save(out_path)
print(f'Saved: {os.path.abspath(out_path)}')
